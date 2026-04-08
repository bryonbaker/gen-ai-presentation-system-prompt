"""
Shared helper library for generating Red Hat branded PowerPoint presentations.

Usage:
    from presentation_helpers import *

    prs = create_presentation(TEMPLATE_PATH)
    slide = add_blank_slide(prs)
    # ... build your slides ...
    prs.save(output_path)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ============================================================
# PATHS
# ============================================================
TEMPLATE_PATH = './FinOS/0 - clean template.pptx'
VENV_PATH = './FinOS/.venv'

# ============================================================
# RED HAT BRAND COLORS
# ============================================================
RH_RED = RGBColor(0xEE, 0x00, 0x00)
RH_RED_DARK = RGBColor(0xCC, 0x00, 0x00)
RH_BLACK = RGBColor(0x15, 0x15, 0x15)
RH_DARK = RGBColor(0x29, 0x29, 0x29)
RH_BODY = RGBColor(0x1F, 0x1F, 0x1F)
RH_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RH_GRAY = RGBColor(0x99, 0x99, 0x99)
RH_GRAY_200 = RGBColor(0xD2, 0xD2, 0xD2)
RH_GRAY_100 = RGBColor(0xED, 0xED, 0xED)
RH_GRAY_50 = RGBColor(0xF5, 0xF5, 0xF5)
RH_BLUE = RGBColor(0x00, 0x66, 0xCC)
RH_GREEN = RGBColor(0x2E, 0x7D, 0x32)
RH_ORANGE = RGBColor(0xC6, 0x51, 0x00)
RH_PURPLE = RGBColor(0x5E, 0x35, 0xB1)
RH_CYAN = RGBColor(0x00, 0x79, 0x6B)

# Light backgrounds (for cards and highlights)
RH_RED_LIGHT = RGBColor(0xFC, 0xEA, 0xE9)
RH_BLUE_LIGHT = RGBColor(0xE3, 0xF2, 0xFD)
RH_GREEN_LIGHT = RGBColor(0xE8, 0xF5, 0xE9)
RH_ORANGE_LIGHT = RGBColor(0xFF, 0xF3, 0xE0)
RH_PURPLE_LIGHT = RGBColor(0xED, 0xE7, 0xF6)

# ============================================================
# FONTS
# ============================================================
FONT_DISPLAY = 'Red Hat Display'
FONT_TEXT = 'Red Hat Text'
FONT_MONO = 'Red Hat Mono'

# ============================================================
# SLIDE DIMENSIONS (EMU)
# ============================================================
SW = 9144000  # 10 inches
SH = 5143500  # 5.625 inches

# ============================================================
# LAYOUT INDEX
# ============================================================
BLANK_LAYOUT_INDEX = 70


# ============================================================
# PRESENTATION CREATION
# ============================================================

def create_presentation(template_path=TEMPLATE_PATH):
    """Create a new presentation from the template with all existing slides removed."""
    prs = Presentation(template_path)
    # Remove existing slides
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])
    return prs


def add_blank_slide(prs):
    """Add a blank slide using the BLANK layout."""
    return prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT_INDEX])


# ============================================================
# TEXT HELPERS
# ============================================================

def add_textbox(slide, left, top, width, height):
    """Add a textbox and return it."""
    return slide.shapes.add_textbox(left, top, width, height)


def set_text(tf, text, font_name=FONT_TEXT, size=Pt(14), color=RH_BODY,
             bold=False, alignment=PP_ALIGN.LEFT):
    """Set text on a text frame (clears existing)."""
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold
    return p


def add_paragraph(tf, text, font_name=FONT_TEXT, size=Pt(14), color=RH_BODY,
                  bold=False, alignment=PP_ALIGN.LEFT, space_before=Pt(4),
                  space_after=Pt(2)):
    """Add a paragraph to a text frame."""
    p = tf.add_paragraph()
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold
    return p


def add_bullet(tf, text, font_name=FONT_TEXT, size=Pt(13), color=RH_BODY,
               bold=False, level=0, space_before=Pt(2)):
    """Add a bullet point paragraph."""
    p = tf.add_paragraph()
    p.level = level
    p.space_before = space_before
    p.space_after = Pt(1)
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold
    return p


def add_bold_then_normal(p, bold_text, normal_text, font_name=FONT_TEXT,
                         size=Pt(13), bold_color=RH_BLACK, normal_color=RH_BODY):
    """Add a run with bold text followed by normal text to a paragraph."""
    r1 = p.add_run()
    r1.text = bold_text
    r1.font.name = font_name
    r1.font.size = size
    r1.font.color.rgb = bold_color
    r1.font.bold = True
    r2 = p.add_run()
    r2.text = normal_text
    r2.font.name = font_name
    r2.font.size = size
    r2.font.color.rgb = normal_color
    r2.font.bold = False
    return p


def set_notes(slide, text):
    """Set speaker notes on a slide."""
    slide.notes_slide.notes_text_frame.text = text


# ============================================================
# SHAPE HELPERS
# ============================================================

def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def add_circle(slide, left, top, size, fill_color):
    """Add a circle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


# ============================================================
# SLIDE ELEMENT HELPERS
# ============================================================

def add_red_accent_bar(slide):
    """Add the red left accent bar (used on content slides)."""
    add_rect(slide, 0, 0, Pt(5), SH, RH_RED)


def add_section_label(slide, text, left=Inches(0.6), top=Inches(0.4)):
    """Add a section label (uppercase, red, small)."""
    tb = add_textbox(slide, left, top, Inches(4), Pt(20))
    set_text(tb.text_frame, text.upper(), FONT_TEXT, Pt(10), RH_RED, bold=True)
    return tb


def add_slide_title(slide, text, left=Inches(0.6), top=Inches(0.7), width=Inches(8.5)):
    """Add a slide title (h2)."""
    tb = add_textbox(slide, left, top, width, Pt(40))
    set_text(tb.text_frame, text, FONT_DISPLAY, Pt(26), RH_BLACK, bold=True)
    return tb


def add_slide_subtitle(slide, text, left=Inches(0.6), top=Inches(1.2), width=Inches(8)):
    """Add a slide subtitle."""
    tb = add_textbox(slide, left, top, width, Pt(24))
    set_text(tb.text_frame, text, FONT_TEXT, Pt(14), RH_DARK)
    return tb


def add_callout(slide, text, left=Inches(0.6), top=Inches(4.2),
                width=Inches(8.8), height=Inches(0.5)):
    """Add a callout box with red left border."""
    add_rect(slide, left, top, Pt(4), height, RH_RED)
    add_rect(slide, left + Pt(5), top, width, height, RH_GRAY_50)
    tb = add_textbox(slide, left + Pt(24), top + Pt(5), width - Pt(34), height - Pt(10))
    tf = tb.text_frame
    tf.word_wrap = True
    set_text(tf, text, FONT_TEXT, Pt(12), RH_BODY)
    return tf


# ============================================================
# COMPOSITE SLIDE BUILDERS
# ============================================================

def make_card(slide, left, top, width, height, title, bullets,
              border_color=RH_GRAY_200, title_color=RH_BLACK, bg_color=RH_WHITE):
    """Create a card with title and bullets."""
    shape = add_rounded_rect(slide, left, top, width, height, bg_color, border_color)
    tb = add_textbox(slide, left + Pt(14), top + Pt(10), width - Pt(28), height - Pt(20))
    tf = tb.text_frame
    tf.word_wrap = True
    set_text(tf, title, FONT_DISPLAY, Pt(14), title_color, bold=True)
    for bullet in bullets:
        add_bullet(tf, bullet, size=Pt(11), space_before=Pt(3))
    return shape


def make_title_splash(prs, title, subtitle, meta_lines=None):
    """Create a red title splash slide (opening or closing)."""
    slide = add_blank_slide(prs)
    add_rect(slide, 0, 0, SW, SH, RH_RED)

    tb = add_textbox(slide, Inches(1), Inches(1.0), Inches(8), Inches(1.2))
    set_text(tb.text_frame, title, FONT_DISPLAY, Pt(40), RH_WHITE,
             bold=True, alignment=PP_ALIGN.CENTER)

    # Divider
    add_rect(slide, Inches(4.2), Inches(2.2), Inches(1.6), Pt(3),
             RGBColor(0xFF, 0x99, 0x99))

    tb = add_textbox(slide, Inches(1.5), Inches(2.5), Inches(7), Inches(0.8))
    set_text(tb.text_frame, subtitle, FONT_TEXT, Pt(18), RH_WHITE,
             alignment=PP_ALIGN.CENTER)

    if meta_lines:
        tb = add_textbox(slide, Inches(1.5), Inches(3.5), Inches(7), Inches(0.8))
        tf = tb.text_frame
        tf.word_wrap = True
        set_text(tf, meta_lines[0], FONT_TEXT, Pt(11),
                 RGBColor(0xFF, 0xCC, 0xCC), alignment=PP_ALIGN.CENTER)
        for line in meta_lines[1:]:
            add_paragraph(tf, line, FONT_TEXT, Pt(11),
                          RGBColor(0xFF, 0xCC, 0xCC), alignment=PP_ALIGN.CENTER)

    return slide


def make_section_slide(prs, number, title, subtitle=""):
    """Create a section transition slide with red left panel."""
    slide = add_blank_slide(prs)
    # Red left panel (40%)
    add_rect(slide, 0, 0, Inches(4), SH, RH_RED)
    # Section number
    tb = add_textbox(slide, Inches(1.2), Inches(1.5), Inches(2), Inches(1.5))
    set_text(tb.text_frame, f"{number:02d}", FONT_DISPLAY, Pt(72),
             RGBColor(0xFF, 0x66, 0x66), bold=True, alignment=PP_ALIGN.RIGHT)
    # Title + subtitle
    tb = add_textbox(slide, Inches(4.5), Inches(1.5), Inches(5), Inches(1.8))
    tf = tb.text_frame
    tf.word_wrap = True
    set_text(tf, title, FONT_DISPLAY, Pt(28), RH_BLACK, bold=True)
    if subtitle:
        add_paragraph(tf, subtitle, FONT_TEXT, Pt(14), RH_DARK, space_before=Pt(12))
    return slide


def make_content_slide(prs, section_label, title, subtitle=""):
    """Create a standard content slide with red accent bar, label, title, and subtitle."""
    slide = add_blank_slide(prs)
    add_red_accent_bar(slide)
    add_section_label(slide, section_label)
    add_slide_title(slide, title)
    if subtitle:
        add_slide_subtitle(slide, subtitle)
    return slide


def make_dark_slide(prs):
    """Create a dark background slide for high-impact messages."""
    slide = add_blank_slide(prs)
    add_rect(slide, 0, 0, SW, SH, RH_BLACK)
    return slide


def make_thank_you(prs, subtitle="Questions & Discussion"):
    """Create a red thank you slide."""
    slide = add_blank_slide(prs)
    add_rect(slide, 0, 0, SW, SH, RH_RED)

    tb = add_textbox(slide, Inches(1), Inches(1.5), Inches(8), Inches(1.2))
    set_text(tb.text_frame, "Thank You", FONT_DISPLAY, Pt(44), RH_WHITE,
             bold=True, alignment=PP_ALIGN.CENTER)

    add_rect(slide, Inches(4.2), Inches(2.7), Inches(1.6), Pt(3),
             RGBColor(0xFF, 0x99, 0x99))

    tb = add_textbox(slide, Inches(1.5), Inches(3.0), Inches(7), Inches(0.6))
    set_text(tb.text_frame, subtitle, FONT_TEXT, Pt(20), RH_WHITE,
             alignment=PP_ALIGN.CENTER)

    return slide
