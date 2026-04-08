#!/usr/bin/env python3
"""Generate Financial AI presentation from Red Hat template."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

TEMPLATE = './FinOS/0 - clean template.pptx'
OUTPUT = './FinOS/Evaluating Financial AI.pptx'

# Red Hat brand colors
RH_RED = RGBColor(0xEE, 0x00, 0x00)
RH_BLACK = RGBColor(0x15, 0x15, 0x15)
RH_DARK = RGBColor(0x29, 0x29, 0x29)
RH_BODY = RGBColor(0x1F, 0x1F, 0x1F)
RH_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RH_GRAY = RGBColor(0x99, 0x99, 0x99)
RH_GRAY_200 = RGBColor(0xD2, 0xD2, 0xD2)
RH_BLUE = RGBColor(0x00, 0x66, 0xCC)
RH_GREEN = RGBColor(0x2E, 0x7D, 0x32)
RH_ORANGE = RGBColor(0xC6, 0x51, 0x00)
RH_PURPLE = RGBColor(0x5E, 0x35, 0xB1)

# Light backgrounds
RH_RED_LIGHT = RGBColor(0xFC, 0xEA, 0xE9)
RH_BLUE_LIGHT = RGBColor(0xE3, 0xF2, 0xFD)
RH_GREEN_LIGHT = RGBColor(0xE8, 0xF5, 0xE9)
RH_ORANGE_LIGHT = RGBColor(0xFF, 0xF3, 0xE0)

FONT_DISPLAY = 'Red Hat Display'
FONT_TEXT = 'Red Hat Text'
FONT_MONO = 'Red Hat Mono'

# Slide dimensions in EMU (10" x 5.625")
SW = 9144000
SH = 5143500


def add_textbox(slide, left, top, width, height):
    """Add a textbox and return it."""
    return slide.shapes.add_textbox(left, top, width, height)


def set_text(tf, text, font_name=FONT_TEXT, size=Pt(14), color=RH_BODY,
             bold=False, alignment=PP_ALIGN.LEFT):
    """Set text on a text frame (clears existing)."""
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold
    return p


def add_paragraph(tf, text, font_name=FONT_TEXT, size=Pt(14), color=RH_BODY,
                  bold=False, alignment=PP_ALIGN.LEFT, space_before=Pt(4),
                  space_after=Pt(2)):
    """Add a paragraph to a text frame."""
    p = tf.add_paragraph()
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold
    return p


def add_bullet(tf, text, font_name=FONT_TEXT, size=Pt(13), color=RH_BODY,
               bold=False, level=0, space_before=Pt(2)):
    """Add a bullet point paragraph."""
    p = tf.add_paragraph()
    p.level = level
    p.space_before = space_before
    p.space_after = Pt(1)
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold
    return p


def add_bold_then_normal(p, bold_text, normal_text, font_name=FONT_TEXT,
                         size=Pt(13), bold_color=RH_BLACK, normal_color=RH_BODY):
    """Add a run with bold text followed by normal text to a paragraph."""
    r1 = p.add_run()
    r1.text = bold_text
    r1.font.name = font_name
    r1.font.size = size
    r1.font.color.rgb = bold_color
    r1.font.bold = True
    r2 = p.add_run()
    r2.text = normal_text
    r2.font.name = font_name
    r2.font.size = size
    r2.font.color.rgb = normal_color
    r2.font.bold = False
    return p


def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def add_red_accent_bar(slide):
    """Add the red left accent bar."""
    add_rect(slide, 0, 0, Pt(5), SH, RH_RED)


def add_section_label(slide, text, left=Inches(0.6), top=Inches(0.4)):
    """Add a section label (uppercase, red, small)."""
    tb = add_textbox(slide, left, top, Inches(4), Pt(20))
    set_text(tb.text_frame, text.upper(), FONT_TEXT, Pt(10), RH_RED, bold=True)
    return tb


def add_slide_title(slide, text, left=Inches(0.6), top=Inches(0.7), width=Inches(8.5)):
    """Add a slide title (h2)."""
    tb = add_textbox(slide, left, top, width, Pt(40))
    set_text(tb.text_frame, text, FONT_DISPLAY, Pt(26), RH_BLACK, bold=True)
    return tb


def add_slide_subtitle(slide, text, left=Inches(0.6), top=Inches(1.2), width=Inches(8)):
    """Add a slide subtitle."""
    tb = add_textbox(slide, left, top, width, Pt(24))
    set_text(tb.text_frame, text, FONT_TEXT, Pt(14), RH_DARK)
    return tb


def make_card(slide, left, top, width, height, title, bullets,
              border_color=RH_GRAY_200, title_color=RH_BLACK, bg_color=RH_WHITE):
    """Create a card with title and bullets."""
    shape = add_rounded_rect(slide, left, top, width, height, bg_color, border_color)
    tb = add_textbox(slide, left + Pt(14), top + Pt(10), width - Pt(28), height - Pt(20))
    tf = tb.text_frame
    tf.word_wrap = True
    set_text(tf, title, FONT_DISPLAY, Pt(14), title_color, bold=True)
    for bullet in bullets:
        add_bullet(tf, bullet, size=Pt(11), space_before=Pt(3))
    return shape


# ============================================================
# BUILD PRESENTATION
# ============================================================

prs = Presentation(TEMPLATE)

# Remove existing slides (keep just the slide masters/layouts)
while len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[0].rId
    prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

BLANK_LAYOUT = prs.slide_layouts[70]  # BLANK
TITLE_LAYOUT = prs.slide_layouts[0]   # TITLE
TITLE1_LAYOUT = prs.slide_layouts[1]  # TITLE_1


# ============================================================
# SLIDE 0: Title splash (red background)
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(slide, 0, 0, SW, SH, RH_RED)

# Title
tb = add_textbox(slide, Inches(1), Inches(1.0), Inches(8), Inches(1.2))
set_text(tb.text_frame, "Evaluating Financial AI", FONT_DISPLAY, Pt(40), RH_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Divider line
add_rect(slide, Inches(4.2), Inches(2.2), Inches(1.6), Pt(3), RGBColor(0xFF, 0x99, 0x99))

# Subtitle
tb = add_textbox(slide, Inches(1.5), Inches(2.5), Inches(7), Inches(0.8))
set_text(tb.text_frame, "A Framework for Benchmarking and Governing\nLLMs in Finance", FONT_TEXT, Pt(18), RH_WHITE, alignment=PP_ALIGN.CENTER)

# Source
tb = add_textbox(slide, Inches(1.5), Inches(3.5), Inches(7), Inches(0.6))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "Based on: Evaluation and Benchmarking Suite for Financial LLMs and Agents", FONT_TEXT, Pt(11), RGBColor(0xFF, 0xCC, 0xCC), alignment=PP_ALIGN.CENTER)
add_paragraph(tf, "SecureFinAI Lab, Columbia University (2026)", FONT_TEXT, Pt(11), RGBColor(0xFF, 0xCC, 0xCC), alignment=PP_ALIGN.CENTER)

slide.notes_slide.notes_text_frame.text = (
    "KEY MESSAGE: Introduce the presentation topic — evaluating and governing financial AI.\n\n"
    "- This presentation is based on a 2026 research paper from Columbia University's SecureFinAI Lab\n"
    "- The paper proposes a comprehensive framework for evaluating financial LLMs and AI agents\n"
    "- Financial AI is moving from experimentation to production — governance is now critical\n"
    "- The framework addresses three concerns: model quality, agent safety, and regulatory compliance\n"
    "- This is relevant to any financial institution deploying or evaluating generative AI\n"
    "- We will walk through the problem, the lifecycle, benchmarking, agents, and governance\n"
    "- The goal is to frame AI as governable financial infrastructure, not experimental technology"
)


# ============================================================
# SLIDE 1: Section — The Problem
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
# Red left panel (40%)
add_rect(slide, 0, 0, Inches(4), SH, RH_RED)
# Section number
tb = add_textbox(slide, Inches(1.2), Inches(1.5), Inches(2), Inches(1.5))
set_text(tb.text_frame, "01", FONT_DISPLAY, Pt(72), RGBColor(0xFF, 0x66, 0x66), bold=True, alignment=PP_ALIGN.RIGHT)
# Section title
tb = add_textbox(slide, Inches(4.5), Inches(1.5), Inches(5), Inches(1.8))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "Why Financial AI Requires a Different Operating Model", FONT_DISPLAY, Pt(28), RH_BLACK, bold=True)
add_paragraph(tf, "Understanding the unique risks of deploying AI in regulated financial environments", FONT_TEXT, Pt(14), RH_DARK, space_before=Pt(12))

slide.notes_slide.notes_text_frame.text = (
    "KEY MESSAGE: Financial AI cannot be treated like general-purpose consumer AI.\n\n"
    "- Finance is a high-stakes, heavily regulated domain — errors have real consequences\n"
    "- General-purpose LLMs were not designed for financial reasoning and frequently make mistakes\n"
    "- We are now seeing autonomous AI agents making financial decisions, not just answering questions\n"
    "- This creates entirely new risk categories that existing IT governance doesn't cover\n"
    "- Regulators are beginning to ask how AI decisions can be traced and audited\n"
    "- The next slide breaks down the specific risk categories we need to address"
)


# ============================================================
# SLIDE 2: Australian Regulatory Context
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_red_accent_bar(slide)
add_section_label(slide, "Regulatory Context")
add_slide_title(slide, "AI governance in the Australian regulatory context")
add_slide_subtitle(slide, "Australia has not introduced dedicated AI regulation for banks. APRA expects existing frameworks to apply.")

# Two-column layout
col_left = Inches(0.6)
col_right = Inches(5.1)
col_w = Inches(4.3)

# Left column: APRA's Position + Standards
# APRA Position card
apra_h = Inches(1.6)
shape = add_rounded_rect(slide, col_left, Inches(1.8), col_w, apra_h, RH_WHITE, RH_RED)
tb = add_textbox(slide, col_left + Pt(14), Inches(1.8) + Pt(10), col_w - Pt(28), apra_h - Pt(20))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "APRA's Position", FONT_DISPLAY, Pt(14), RH_RED, bold=True)
add_paragraph(tf, '"We remain of the view that our existing regulatory framework is sufficient to capture the use of AI."', FONT_TEXT, Pt(10), RH_BODY, space_before=Pt(6))
add_paragraph(tf, "AI must be managed within existing frameworks for:", FONT_TEXT, Pt(11), RH_BODY, space_before=Pt(6))
add_bullet(tf, "Model risk management", size=Pt(11), space_before=Pt(2))
add_bullet(tf, "Operational risk", size=Pt(11))
add_bullet(tf, "Technology risk and data security", size=Pt(11))

# Relevant Standards card
std_top = Inches(3.5)
std_h = Inches(1.6)
shape = add_rounded_rect(slide, col_left, std_top, col_w, std_h, RGBColor(0xF5, 0xF5, 0xF5), RH_GRAY_200)
tb = add_textbox(slide, col_left + Pt(14), std_top + Pt(10), col_w - Pt(28), std_h - Pt(20))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "Relevant APRA Standards", FONT_DISPLAY, Pt(13), RH_BLACK, bold=True)
p = add_bullet(tf, "", size=Pt(11), space_before=Pt(6))
p.clear()
add_bold_then_normal(p, "CPS 220 ", "\u2014 Risk Management: governance of material risks", size=Pt(11))
p = add_bullet(tf, "", size=Pt(11))
p.clear()
add_bold_then_normal(p, "CPS 231 ", "\u2014 Outsourcing: third-party AI services and model providers", size=Pt(11))
p = add_bullet(tf, "", size=Pt(11))
p.clear()
add_bold_then_normal(p, "CPS 234 ", "\u2014 Information Security: data protection and resilience", size=Pt(11))
p = add_bullet(tf, "", size=Pt(11))
p.clear()
add_bold_then_normal(p, "Model Risk ", "\u2014 Validation and monitoring of analytical models", size=Pt(11))

# Right column: Industry Context + Why This Matters
# Industry Context card
ctx_h = Inches(1.6)
shape = add_rounded_rect(slide, col_right, Inches(1.8), col_w, ctx_h, RH_ORANGE_LIGHT, RH_ORANGE)
tb = add_textbox(slide, col_right + Pt(14), Inches(1.8) + Pt(10), col_w - Pt(28), ctx_h - Pt(20))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "Other Regulators Are Moving Faster", FONT_DISPLAY, Pt(14), RH_ORANGE, bold=True)
add_paragraph(tf, "Singapore MAS AI Risk Management Framework (2025):", FONT_TEXT, Pt(11), RH_BODY, space_before=Pt(6))
add_bullet(tf, "Explicit guidance for AI model governance", size=Pt(11), space_before=Pt(2))
add_bullet(tf, "Explainability and fairness requirements", size=Pt(11))
add_bullet(tf, "Lifecycle monitoring expectations", size=Pt(11))
add_paragraph(tf, "Australia is currently taking a technology-neutral approach.", FONT_TEXT, Pt(11), RH_BODY, bold=True, space_before=Pt(6))

# Why This Matters card
why_top = Inches(3.5)
why_h = Inches(1.6)
shape = add_rounded_rect(slide, col_right, why_top, col_w, why_h, RH_GREEN_LIGHT, RH_GREEN)
tb = add_textbox(slide, col_right + Pt(14), why_top + Pt(10), col_w - Pt(28), why_h - Pt(20))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "Why This Presentation Matters", FONT_DISPLAY, Pt(14), RH_GREEN, bold=True)
add_paragraph(tf, "If APRA expects AI to be governed under existing frameworks, banks need practical ways to implement:", FONT_TEXT, Pt(11), RH_BODY, space_before=Pt(6))
add_bullet(tf, "Model validation for AI systems", size=Pt(11), space_before=Pt(2))
add_bullet(tf, "Monitoring of autonomous AI agents", size=Pt(11))
add_bullet(tf, "Transparency and auditability of AI models", size=Pt(11))

slide.notes_slide.notes_text_frame.text = (
    "KEY MESSAGE: APRA has not created new AI regulations — it expects banks to govern AI under existing prudential standards. This creates both a challenge and an opportunity.\n\n"
    "- APRA's Chris Gower stated explicitly that the existing regulatory framework is considered sufficient to capture AI use by banks, insurers, and super funds\n"
    "- This means AI deployments must comply with CPS 220 (risk management), CPS 231 (outsourcing for third-party models), CPS 234 (information security), and existing model risk management practices\n"
    "- In contrast, Singapore's MAS published explicit AI risk management guidelines in 2025, covering model governance, explainability, fairness, and lifecycle monitoring\n"
    "- Australia's technology-neutral approach means banks must interpret how existing standards apply to AI — the frameworks in this presentation provide practical implementation guidance\n"
    "- The FinOS governance architecture maps directly onto APRA's existing requirements: model evaluation satisfies model risk management, agent monitoring satisfies operational risk, and transparency satisfies CPS 231 vendor obligations\n"
    "- For Australian banks, this is not about waiting for new regulation — it is about applying existing obligations to a new technology class\n"
    "- Proactive governance now positions the bank ahead of potential future regulatory tightening"
)


# ============================================================
# SLIDE 3: Risk Categories
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_red_accent_bar(slide)
add_section_label(slide, "The Problem")
add_slide_title(slide, "Financial AI introduces new risk categories")
add_slide_subtitle(slide, "Financial AI must be evaluated differently from consumer AI.")

# Card 1: Model Risk (top-left)
card_w = Inches(4.3)
card_h = Inches(1.7)
card_left1 = Inches(0.6)
card_left2 = Inches(5.1)
card_top1 = Inches(1.8)
card_top2 = Inches(3.6)

make_card(slide, card_left1, card_top1, card_w, card_h,
          "Model Risk (Financial Reasoning)",
          ["Interpreting earnings reports and SEC filings",
           "Analysing financial statements",
           "Numerical reasoning in investment analysis",
           "Errors affect risk models and trading decisions"],
          border_color=RH_RED, title_color=RH_RED)

# Card 2: Operational Risk (top-right)
make_card(slide, card_left2, card_top1, card_w, card_h,
          "Operational Risk (Autonomous Agents)",
          ["Retrieving market data autonomously",
           "Analysing financial news",
           "Generating investment insights",
           "Non-deterministic behaviour requires monitoring"],
          border_color=RH_ORANGE, title_color=RH_ORANGE)

# Card 3: Regulatory Risk (bottom-left)
make_card(slide, card_left1, card_top2, card_w, Inches(1.4),
          "Regulatory & Compliance Risk",
          ["Training data provenance is unclear",
           "Licensing restrictions are unknown",
           "Outputs cannot be audited by regulators"],
          border_color=RH_PURPLE, title_color=RH_PURPLE)

# Card 4: Governance Gap (bottom-right)
make_card(slide, card_left2, card_top2, card_w, Inches(1.4),
          "Governance Gap",
          ["No standard for benchmarking financial AI",
           "No framework for evaluating agents in production",
           "No widely adopted transparency standard"],
          border_color=RH_GRAY_200, title_color=RH_BLACK)

slide.notes_slide.notes_text_frame.text = (
    "KEY MESSAGE: Financial AI introduces four distinct risk categories that require new governance approaches.\n\n"
    "- Model Risk: LLMs can misinterpret financial statements, produce incorrect numerical analysis, and hallucinate in regulatory documents — this directly impacts trading and risk decisions\n"
    "- Operational Risk: AI agents operate autonomously — retrieving data, analysing news, generating insights — and their non-deterministic behaviour makes them unpredictable\n"
    "- Regulatory Risk: Many models have unclear training data provenance and licensing — regulators may require auditability that current models cannot provide\n"
    "- Governance Gap: There is currently no industry standard for benchmarking financial AI or evaluating agents in production\n"
    "- These four risks together mean we need a fundamentally different approach to AI governance in finance\n"
    "- The paper we are reviewing proposes a framework to address all four"
)


# ============================================================
# SLIDE 4: Core Challenge
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_red_accent_bar(slide)
add_section_label(slide, "The Problem")
add_slide_title(slide, "The core challenge")
add_slide_subtitle(slide, "Financial institutions deploying GenAI face three critical questions:")

# Three question cards
q_w = Inches(2.9)
q_h = Inches(2.0)
q_top = Inches(1.9)
q_gap = Inches(0.15)
q_left1 = Inches(0.6)
q_left2 = q_left1 + q_w + q_gap
q_left3 = q_left2 + q_w + q_gap

# Card 1
shape = add_rounded_rect(slide, q_left1, q_top, q_w, q_h, RH_RED_LIGHT, RH_RED)
tb = add_textbox(slide, q_left1 + Pt(14), q_top + Pt(14), q_w - Pt(28), q_h - Pt(28))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "Which models can we trust?", FONT_DISPLAY, Pt(16), RH_BLACK, bold=True, alignment=PP_ALIGN.CENTER)
add_paragraph(tf, "How do we evaluate financial reasoning capability and accuracy?", FONT_TEXT, Pt(12), RH_BODY, alignment=PP_ALIGN.CENTER, space_before=Pt(10))

# Card 2
shape = add_rounded_rect(slide, q_left2, q_top, q_w, q_h, RH_ORANGE_LIGHT, RH_ORANGE)
tb = add_textbox(slide, q_left2 + Pt(14), q_top + Pt(14), q_w - Pt(28), q_h - Pt(28))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "How do we safely operate AI agents?", FONT_DISPLAY, Pt(16), RH_BLACK, bold=True, alignment=PP_ALIGN.CENTER)
add_paragraph(tf, "How do we monitor autonomous systems in production?", FONT_TEXT, Pt(12), RH_BODY, alignment=PP_ALIGN.CENTER, space_before=Pt(10))

# Card 3
shape = add_rounded_rect(slide, q_left3, q_top, q_w, q_h, RH_GREEN_LIGHT, RH_GREEN)
tb = add_textbox(slide, q_left3 + Pt(14), q_top + Pt(14), q_w - Pt(28), q_h - Pt(28))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "How do we meet audit requirements?", FONT_DISPLAY, Pt(16), RH_BLACK, bold=True, alignment=PP_ALIGN.CENTER)
add_paragraph(tf, "How do we ensure regulatory compliance and model transparency?", FONT_TEXT, Pt(12), RH_BODY, alignment=PP_ALIGN.CENTER, space_before=Pt(10))

# Callout
add_rect(slide, Inches(0.6), Inches(4.2), Pt(4), Inches(0.5), RH_RED)
shape = add_rect(slide, Inches(0.65), Inches(4.2), Inches(8.8), Inches(0.5), RGBColor(0xF5, 0xF5, 0xF5))
tb = add_textbox(slide, Inches(0.9), Inches(4.25), Inches(8.4), Inches(0.4))
set_text(tb.text_frame, "This presentation outlines a framework addressing all three.", FONT_TEXT, Pt(13), RH_BODY)

slide.notes_slide.notes_text_frame.text = (
    "KEY MESSAGE: Every bank deploying GenAI must answer three fundamental questions.\n\n"
    "- Question 1 — Trust: How do we know a model is accurate enough for financial tasks? General benchmarks don't test financial reasoning\n"
    "- Question 2 — Operations: AI agents are making autonomous decisions in production — how do we monitor and control them?\n"
    "- Question 3 — Compliance: Regulators will ask how AI decisions were made and whether the models meet transparency requirements\n"
    "- These are not theoretical concerns — Google's AI Overview produced 43% inaccurate finance summaries, and a ChatGPT security incident exposed user financial data\n"
    "- The rest of this presentation walks through a framework that addresses all three questions\n"
    "- Think of this as: model evaluation, agent monitoring, and model transparency — three layers of governance"
)


# ============================================================
# SLIDE 4: Section — Lifecycle
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(slide, 0, 0, Inches(4), SH, RH_RED)
tb = add_textbox(slide, Inches(1.2), Inches(1.5), Inches(2), Inches(1.5))
set_text(tb.text_frame, "02", FONT_DISPLAY, Pt(72), RGBColor(0xFF, 0x66, 0x66), bold=True, alignment=PP_ALIGN.RIGHT)
tb = add_textbox(slide, Inches(4.5), Inches(1.5), Inches(5), Inches(1.8))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "The Financial AI Lifecycle", FONT_DISPLAY, Pt(28), RH_BLACK, bold=True)
add_paragraph(tf, "From experimentation to governed deployment", FONT_TEXT, Pt(14), RH_DARK, space_before=Pt(12))

slide.notes_slide.notes_text_frame.text = (
    "KEY MESSAGE: Transition to the lifecycle section — how financial AI has matured over three years.\n\n"
    "- The financial AI field has evolved rapidly from 2023 to 2025\n"
    "- This section maps that evolution into three distinct maturity stages\n"
    "- Understanding where the industry is helps us understand what governance is needed now\n"
    "- Most banks are currently between Readiness and Governance — deploying models but lacking governance frameworks\n"
    "- The research community has been building tools and frameworks at each stage\n"
    "- The next slide details each stage and its key outputs"
)


# ============================================================
# SLIDE 5: Three Stages
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_red_accent_bar(slide)
add_section_label(slide, "Lifecycle")
add_slide_title(slide, "Three stages of financial LLM maturity")
add_slide_subtitle(slide, "Financial AI is evolving through three maturity stages.")

# Three stage boxes in a row with arrows
stg_w = Inches(2.7)
stg_h = Inches(2.6)
stg_top = Inches(1.8)
stg_left1 = Inches(0.6)
stg_left2 = Inches(3.55)
stg_left3 = Inches(6.5)

# Stage 1: Exploration
shape = add_rounded_rect(slide, stg_left1, stg_top, stg_w, stg_h, RH_WHITE, RH_GRAY_200)
tb = add_textbox(slide, stg_left1 + Pt(12), stg_top + Pt(10), stg_w - Pt(24), stg_h - Pt(20))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "2023", FONT_MONO, Pt(11), RH_ORANGE, bold=True)
add_paragraph(tf, "Exploration", FONT_DISPLAY, Pt(16), RH_BLACK, bold=True, space_before=Pt(4))
add_paragraph(tf, "Early financial LLM development.", FONT_TEXT, Pt(11), RH_BODY, space_before=Pt(8))
add_bullet(tf, "BloombergGPT, FinGPT", size=Pt(11))
add_bullet(tf, "Can LLMs understand financial language?", size=Pt(11))
add_bullet(tf, "Open FinLLM Leaderboard", size=Pt(11))

# Arrow 1
tb = add_textbox(slide, Inches(3.3), Inches(2.8), Inches(0.3), Inches(0.3))
set_text(tb.text_frame, "\u2192", FONT_TEXT, Pt(20), RH_DARK, alignment=PP_ALIGN.CENTER)

# Stage 2: Readiness
shape = add_rounded_rect(slide, stg_left2, stg_top, stg_w, stg_h, RH_WHITE, RH_GRAY_200)
tb = add_textbox(slide, stg_left2 + Pt(12), stg_top + Pt(10), stg_w - Pt(24), stg_h - Pt(20))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "2024", FONT_MONO, Pt(11), RH_ORANGE, bold=True)
add_paragraph(tf, "Readiness", FONT_DISPLAY, Pt(16), RH_BLACK, bold=True, space_before=Pt(4))
add_paragraph(tf, "Shift to deployment readiness.", FONT_TEXT, Pt(11), RH_BODY, space_before=Pt(8))
add_bullet(tf, "Systematic benchmarking", size=Pt(11))
add_bullet(tf, "Financial datasets & evaluation", size=Pt(11))
add_bullet(tf, "Emergence of financial agents", size=Pt(11))

# Arrow 2
tb = add_textbox(slide, Inches(6.25), Inches(2.8), Inches(0.3), Inches(0.3))
set_text(tb.text_frame, "\u2192", FONT_TEXT, Pt(20), RH_DARK, alignment=PP_ALIGN.CENTER)

# Stage 3: Governance (highlighted)
shape = add_rounded_rect(slide, stg_left3, stg_top, stg_w, stg_h, RH_RED_LIGHT, RH_RED)
tb = add_textbox(slide, stg_left3 + Pt(12), stg_top + Pt(10), stg_w - Pt(24), stg_h - Pt(20))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "2025", FONT_MONO, Pt(11), RH_RED, bold=True)
add_paragraph(tf, "Governance", FONT_DISPLAY, Pt(16), RH_RED, bold=True, space_before=Pt(4))
add_paragraph(tf, "Safe and compliant deployment.", FONT_TEXT, Pt(11), RH_BODY, space_before=Pt(8))
add_bullet(tf, "AI governance frameworks", size=Pt(11))
add_bullet(tf, "Operational evaluation of agents", size=Pt(11))
add_bullet(tf, "Model transparency standards", size=Pt(11))

slide.notes_slide.notes_text_frame.text = (
    "KEY MESSAGE: Financial AI has progressed through three maturity stages — we are now in the governance era.\n\n"
    "- 2023 Exploration: Early models like BloombergGPT (50B parameters, trained on financial data) and open-source FinGPT proved LLMs could understand financial language\n"
    "- The Open FinLLM Leaderboard was created during this stage to begin standardising evaluation\n"
    "- 2024 Readiness: The focus shifted from 'can it work?' to 'is it production-ready?' — comprehensive benchmarks and financial agents emerged\n"
    "- 2025 Governance: The current focus is on safety, compliance, hallucination detection, and operational monitoring of deployed systems\n"
    "- This mirrors the DevOps to MLOps evolution — now we need AgentOps for AI agents\n"
    "- The highlighted Governance stage is where most banks need to invest today\n"
    "- Each stage produced specific tools: Leaderboard, evaluation pipelines, and governance frameworks"
)


# ============================================================
# SLIDE 6: Section — Benchmarking
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(slide, 0, 0, Inches(4), SH, RH_RED)
tb = add_textbox(slide, Inches(1.2), Inches(1.5), Inches(2), Inches(1.5))
set_text(tb.text_frame, "03", FONT_DISPLAY, Pt(72), RGBColor(0xFF, 0x66, 0x66), bold=True, alignment=PP_ALIGN.RIGHT)
tb = add_textbox(slide, Inches(4.5), Inches(1.5), Inches(5), Inches(1.8))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "Benchmarking Financial AI", FONT_DISPLAY, Pt(28), RH_BLACK, bold=True)
add_paragraph(tf, "Standardised evaluation across financial tasks", FONT_TEXT, Pt(14), RH_DARK, space_before=Pt(12))

slide.notes_slide.notes_text_frame.text = (
    "KEY MESSAGE: Transition to benchmarking — how we measure whether a financial AI model is fit for purpose.\n\n"
    "- Benchmarking is the first layer of governance — you cannot govern what you cannot measure\n"
    "- Currently, different institutions evaluate models using different criteria, making comparison impossible\n"
    "- The Open FinLLM Leaderboard aims to be a de facto standard, similar to how stress tests standardised risk evaluation\n"
    "- This is a collaborative effort between academia (Columbia), open source (Linux Foundation, PyTorch), and industry\n"
    "- The next slide shows what the leaderboard actually tests"
)


# ============================================================
# SLIDE 7: Leaderboard
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_red_accent_bar(slide)
add_section_label(slide, "Benchmarking")
add_slide_title(slide, "The Open FinLLM Leaderboard")
add_slide_subtitle(slide, "Standardised evaluation across 42 financial datasets, organised into three capability areas.")

# Three capability cards
cap_w = Inches(2.9)
cap_h = Inches(2.2)
cap_top = Inches(1.8)
cap_left1 = Inches(0.6)
cap_left2 = cap_left1 + cap_w + Inches(0.15)
cap_left3 = cap_left2 + cap_w + Inches(0.15)

make_card(slide, cap_left1, cap_top, cap_w, cap_h,
          "Understanding Financial Language",
          ["Information extraction",
           "Sentiment analysis",
           "Classification"],
          border_color=RH_RED, title_color=RH_RED)

make_card(slide, cap_left2, cap_top, cap_w, cap_h,
          "Financial Reasoning",
          ["Answering questions about financial data",
           "Summarising earnings reports",
           "Analysing regulatory filings"],
          border_color=RH_ORANGE, title_color=RH_ORANGE)

make_card(slide, cap_left3, cap_top, cap_w, cap_h,
          "Financial Decision Tasks",
          ["Credit risk scoring",
           "Market forecasting",
           "Trading and investment simulations"],
          border_color=RH_GREEN, title_color=RH_GREEN)

# Callout
add_rect(slide, Inches(0.6), Inches(4.2), Pt(4), Inches(0.7), RH_RED)
shape = add_rect(slide, Inches(0.65), Inches(4.2), Inches(8.8), Inches(0.7), RGBColor(0xF5, 0xF5, 0xF5))
tb = add_textbox(slide, Inches(0.9), Inches(4.25), Inches(8.4), Inches(0.6))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "Hosted on HuggingFace in collaboration with Linux Foundation, PyTorch Foundation, and SecureFinAI Lab. Models ranked by both performance and openness classification.", FONT_TEXT, Pt(12), RH_BODY)

slide.notes_slide.notes_text_frame.text = (
    "KEY MESSAGE: The Open FinLLM Leaderboard tests models across 42 financial datasets in three capability areas.\n\n"
    "- Understanding Financial Language: Can the model extract entities, classify sentiment, and understand financial terminology?\n"
    "- Financial Reasoning: Can it answer complex questions about financial data, summarise 10-K filings, and analyse earnings reports?\n"
    "- Financial Decision Tasks: Can it score credit risk, predict stock movements, and simulate trading decisions?\n"
    "- Models are ranked not just by accuracy but also by their openness classification — transparency matters alongside performance\n"
    "- This is hosted on HuggingFace and is an open community resource — anyone can submit models for evaluation\n"
    "- For a bank, this provides a standardised way to compare models before procurement or deployment\n"
    "- Think of this as a financial stress test for AI models"
)


# ============================================================
# SLIDE 8: Section — Agents
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(slide, 0, 0, Inches(4), SH, RH_RED)
tb = add_textbox(slide, Inches(1.2), Inches(1.5), Inches(2), Inches(1.5))
set_text(tb.text_frame, "04", FONT_DISPLAY, Pt(72), RGBColor(0xFF, 0x66, 0x66), bold=True, alignment=PP_ALIGN.RIGHT)
tb = add_textbox(slide, Inches(4.5), Inches(1.5), Inches(5), Inches(1.8))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "From Models to Agents", FONT_DISPLAY, Pt(28), RH_BLACK, bold=True)
add_paragraph(tf, "Financial AI agents introduce new capabilities and new risks", FONT_TEXT, Pt(14), RH_DARK, space_before=Pt(12))

slide.notes_slide.notes_text_frame.text = (
    "KEY MESSAGE: Transition to agents — the industry is moving beyond static models to autonomous AI systems.\n\n"
    "- Most enterprises are now deploying AI agents, not just LLMs — agents can take actions, use tools, and make decisions\n"
    "- This creates fundamentally different risks compared to a model that simply answers questions\n"
    "- Agents are non-deterministic — the same input may produce different actions each time\n"
    "- Traditional model evaluation is insufficient for agents — we need to evaluate behaviour, not just outputs\n"
    "- The next slide shows example financial agents and why they require new governance approaches"
)


# ============================================================
# SLIDE 9: Agents
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_red_accent_bar(slide)
add_section_label(slide, "Financial Agents")
add_slide_title(slide, "Agents change the risk profile")
add_slide_subtitle(slide, "Financial AI is moving from static models to autonomous agents.")

# Left column: capabilities + callout
tb = add_textbox(slide, Inches(0.6), Inches(1.8), Inches(4.5), Inches(2.5))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "New agent capabilities", FONT_DISPLAY, Pt(15), RH_BLACK, bold=True)
p = add_bullet(tf, "", size=Pt(12), space_before=Pt(8))
p.clear()
add_bold_then_normal(p, "Multi-step reasoning ", "across financial data", size=Pt(12))
p = add_bullet(tf, "", size=Pt(12))
p.clear()
add_bold_then_normal(p, "Tool usage ", "for data retrieval and analysis", size=Pt(12))
p = add_bullet(tf, "", size=Pt(12))
p.clear()
add_bold_then_normal(p, "Autonomous decision-making", "", size=Pt(12))
p = add_bullet(tf, "", size=Pt(12))
p.clear()
add_bold_then_normal(p, "Real-time financial data ", "retrieval", size=Pt(12))

# Callout box
add_rect(slide, Inches(0.6), Inches(3.7), Pt(4), Inches(0.7), RH_RED)
shape = add_rect(slide, Inches(0.65), Inches(3.7), Inches(4.4), Inches(0.7), RGBColor(0xF5, 0xF5, 0xF5))
tb = add_textbox(slide, Inches(0.9), Inches(3.78), Inches(4), Inches(0.55))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "Traditional LLM evaluation measures answers.", FONT_TEXT, Pt(12), RH_BODY, bold=True)
add_paragraph(tf, "Agent evaluation must measure behaviour.", FONT_TEXT, Pt(12), RH_RED, bold=True, space_before=Pt(2))

# Right column: agent cards
agent_left = Inches(5.4)
agent_w = Inches(4.2)
agent_h = Inches(1.0)

make_card(slide, agent_left, Inches(1.8), agent_w, agent_h,
          "Agentic FinSearch",
          ["Autonomous research agent retrieving financial data and analysing news sentiment"],
          bg_color=RH_RED_LIGHT, border_color=RGBColor(0xF5, 0xC3, 0xC0))

make_card(slide, agent_left, Inches(2.9), agent_w, agent_h,
          "Financial Tutor",
          ["AI-powered financial education: CFA prep and financial training at scale"],
          bg_color=RH_ORANGE_LIGHT, border_color=RGBColor(0xE0, 0xC5, 0xA0))

make_card(slide, agent_left, Inches(4.0), agent_w, agent_h,
          "FinSight Agent",
          ["Multi-agent earnings call analysis with self-evaluation and guardrails"],
          bg_color=RH_GREEN_LIGHT, border_color=RGBColor(0xC3, 0xDF, 0xC0))

slide.notes_slide.notes_text_frame.text = (
    "KEY MESSAGE: Financial AI agents introduce multi-step reasoning and autonomous behaviour — evaluation must shift from measuring answers to measuring behaviour.\n\n"
    "- Agentic FinSearch: Autonomously scrapes financial websites, retrieves data, and produces market insights — like having an analyst that works 24/7\n"
    "- Financial Tutor: Demonstrates scalable AI-powered financial education — CFA prep, credit risk training, accessible to millions simultaneously\n"
    "- FinSight Agent: A multi-agent system that analyses earnings calls using coordinated expert agents with built-in self-evaluation and guardrails\n"
    "- The key insight is in the callout: traditional evaluation tests whether an answer is correct — agent evaluation must test whether the behaviour is safe and appropriate\n"
    "- This means we need to trace decisions, monitor tool usage, and evaluate reasoning trajectories — not just final outputs\n"
    "- This is analogous to the difference between testing a calculator versus auditing a trader"
)


# ============================================================
# SLIDE 10: Section — Governance
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(slide, 0, 0, Inches(4), SH, RH_RED)
tb = add_textbox(slide, Inches(1.2), Inches(1.5), Inches(2), Inches(1.5))
set_text(tb.text_frame, "05", FONT_DISPLAY, Pt(72), RGBColor(0xFF, 0x66, 0x66), bold=True, alignment=PP_ALIGN.RIGHT)
tb = add_textbox(slide, Inches(4.5), Inches(1.5), Inches(5), Inches(1.8))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "Governing Financial AI Systems", FONT_DISPLAY, Pt(28), RH_BLACK, bold=True)
add_paragraph(tf, "Operational evaluation and model transparency", FONT_TEXT, Pt(14), RH_DARK, space_before=Pt(12))

slide.notes_slide.notes_text_frame.text = (
    "KEY MESSAGE: Transition to governance — the tools and frameworks for safely operating financial AI.\n\n"
    "- We have covered what financial AI can do and the risks it introduces — now we address how to govern it\n"
    "- Governance in this context means two things: operational monitoring of agents and transparency of models\n"
    "- This section introduces AgentOps (analogous to DevOps/MLOps) and the Model Openness Framework\n"
    "- These are not theoretical proposals — they are being developed by Linux Foundation, Red Hat, and industry partners\n"
    "- The next slides detail the inner/outer loop evaluation model and the model transparency classification"
)


# ============================================================
# SLIDE 11: AgentOps
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_red_accent_bar(slide)
add_section_label(slide, "Governance")
add_slide_title(slide, "AgentOps: Evaluating AI agents in practice")
add_slide_subtitle(slide, "Two layers of governance for the AI agent lifecycle.")

# Two governance cards side by side
gov_w = Inches(4.3)
gov_h = Inches(2.7)
gov_top = Inches(1.8)

# Inner Loop card
shape = add_rounded_rect(slide, Inches(0.6), gov_top, gov_w, gov_h, RH_WHITE, RH_RED)
# Number circle
circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.85), gov_top + Pt(14), Pt(32), Pt(32))
circle.fill.solid()
circle.fill.fore_color.rgb = RH_RED
circle.line.fill.background()
tb_num = add_textbox(slide, Inches(0.85), gov_top + Pt(14), Pt(32), Pt(32))
set_text(tb_num.text_frame, "1", FONT_DISPLAY, Pt(16), RH_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

tb = add_textbox(slide, Inches(1.35), gov_top + Pt(14), Inches(3.3), Pt(28))
set_text(tb.text_frame, "Development Governance (Inner Loop)", FONT_DISPLAY, Pt(13), RH_BLACK, bold=True)

tb = add_textbox(slide, Inches(0.85), gov_top + Pt(54), gov_w - Pt(28), gov_h - Pt(64))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, '"Glass box" evaluation of agent behaviour during development.', FONT_TEXT, Pt(12), RH_BODY)
add_bullet(tf, "Reasoning validation", size=Pt(12), space_before=Pt(8))
add_bullet(tf, "Tool usage testing", size=Pt(12))
add_bullet(tf, "Context adherence", size=Pt(12))
add_bullet(tf, "Safety guardrails", size=Pt(12))

# Outer Loop card
shape = add_rounded_rect(slide, Inches(5.1), gov_top, gov_w, gov_h, RH_WHITE, RH_ORANGE)
circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.35), gov_top + Pt(14), Pt(32), Pt(32))
circle.fill.solid()
circle.fill.fore_color.rgb = RH_ORANGE
circle.line.fill.background()
tb_num = add_textbox(slide, Inches(5.35), gov_top + Pt(14), Pt(32), Pt(32))
set_text(tb_num.text_frame, "2", FONT_DISPLAY, Pt(16), RH_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

tb = add_textbox(slide, Inches(5.85), gov_top + Pt(14), Inches(3.3), Pt(28))
set_text(tb.text_frame, "Operational Governance (Outer Loop)", FONT_DISPLAY, Pt(13), RH_BLACK, bold=True)

tb = add_textbox(slide, Inches(5.35), gov_top + Pt(54), gov_w - Pt(28), gov_h - Pt(64))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "Continuous evaluation of AI systems in production.", FONT_TEXT, Pt(12), RH_BODY)
add_bullet(tf, "Decision tracing", size=Pt(12), space_before=Pt(8))
add_bullet(tf, "Audit logs", size=Pt(12))
add_bullet(tf, "LLM-as-a-Judge safety scoring", size=Pt(12))
add_bullet(tf, "Human-in-the-loop escalation", size=Pt(12))

# Callout
add_rect(slide, Inches(0.6), Inches(4.7), Pt(4), Inches(0.45), RH_RED)
shape = add_rect(slide, Inches(0.65), Inches(4.7), Inches(8.8), Inches(0.45), RGBColor(0xF5, 0xF5, 0xF5))
tb = add_textbox(slide, Inches(0.9), Inches(4.73), Inches(8.4), Inches(0.4))
set_text(tb.text_frame, "Production failures are fed back into training datasets to continuously improve reliability.", FONT_TEXT, Pt(12), RH_BODY)

slide.notes_slide.notes_text_frame.text = (
    "KEY MESSAGE: AgentOps provides two governance loops — development-time testing and production-time monitoring.\n\n"
    "- Inner Loop (Development): This is glass-box evaluation — you can see inside the agent's reasoning process during development and testing\n"
    "- It evaluates: Does the agent select the right tools? Does it follow logical reasoning steps? Does it stay within its context? Do safety guardrails trigger correctly?\n"
    "- Outer Loop (Production): Once deployed, you cannot see inside the agent — you monitor its behaviour through audit trails, decision logs, and safety scoring\n"
    "- LLM-as-a-Judge: Uses a separate AI model to score the deployed agent's outputs for safety, logic, and ethics violations\n"
    "- Human-in-the-loop: High-risk or low-confidence decisions are escalated to human reviewers\n"
    "- Critical feedback loop: Production failures are captured and fed back into training data, continuously improving the system\n"
    "- This is the AgentOps equivalent of CI/CD for DevOps — a closed-loop operational governance model"
)


# ============================================================
# SLIDE 12: Model Openness Framework
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_red_accent_bar(slide)
add_section_label(slide, "Governance")
add_slide_title(slide, "Model Openness Framework")
add_slide_subtitle(slide, 'Avoiding "open-washing" in AI model procurement.')

# Intro text
tb = add_textbox(slide, Inches(0.6), Inches(1.6), Inches(8.5), Pt(22))
set_text(tb.text_frame, "Many models claim to be open but hide training datasets, licensing terms, and development processes.", FONT_TEXT, Pt(13), RH_BODY)

# Table-like layout
tbl_top = Inches(2.1)
tbl_left = Inches(0.6)
tbl_w = Inches(8.8)

# Header row
add_rect(slide, tbl_left, tbl_top, tbl_w, Pt(30), RGBColor(0xF5, 0xF5, 0xF5), RH_RED)
tb = add_textbox(slide, tbl_left + Pt(10), tbl_top + Pt(4), Inches(1.2), Pt(22))
set_text(tb.text_frame, "CLASS", FONT_TEXT, Pt(10), RH_DARK, bold=True)
tb = add_textbox(slide, tbl_left + Inches(1.4), tbl_top + Pt(4), Inches(2), Pt(22))
set_text(tb.text_frame, "LEVEL OF OPENNESS", FONT_TEXT, Pt(10), RH_DARK, bold=True)
tb = add_textbox(slide, tbl_left + Inches(3.5), tbl_top + Pt(4), Inches(5), Pt(22))
set_text(tb.text_frame, "WHAT'S SHARED", FONT_TEXT, Pt(10), RH_DARK, bold=True)

# Row 1: Class III
row_h = Pt(32)
row_top = tbl_top + Pt(32)
add_rect(slide, tbl_left, row_top, tbl_w, row_h, RH_WHITE, RH_GRAY_200)
tb = add_textbox(slide, tbl_left + Pt(10), row_top + Pt(6), Inches(1.2), Pt(20))
set_text(tb.text_frame, "Class III", FONT_TEXT, Pt(12), RH_BLACK, bold=True)
tb = add_textbox(slide, tbl_left + Inches(1.4), row_top + Pt(6), Inches(2), Pt(20))
set_text(tb.text_frame, "Open Model", FONT_TEXT, Pt(12), RH_BODY)
tb = add_textbox(slide, tbl_left + Inches(3.5), row_top + Pt(6), Inches(5), Pt(20))
set_text(tb.text_frame, "Partially open: model weights, card, evaluation results", FONT_TEXT, Pt(12), RH_BODY)

# Row 2: Class II
row_top = row_top + row_h + Pt(2)
add_rect(slide, tbl_left, row_top, tbl_w, row_h, RH_WHITE, RH_GRAY_200)
tb = add_textbox(slide, tbl_left + Pt(10), row_top + Pt(6), Inches(1.2), Pt(20))
set_text(tb.text_frame, "Class II", FONT_TEXT, Pt(12), RH_BLACK, bold=True)
tb = add_textbox(slide, tbl_left + Inches(1.4), row_top + Pt(6), Inches(2), Pt(20))
set_text(tb.text_frame, "Open Tooling", FONT_TEXT, Pt(12), RH_BODY)
tb = add_textbox(slide, tbl_left + Inches(3.5), row_top + Pt(6), Inches(5), Pt(20))
set_text(tb.text_frame, "Training code, evaluation code, inference code, and tools", FONT_TEXT, Pt(12), RH_BODY)

# Row 3: Class I (highlighted)
row_top = row_top + row_h + Pt(2)
add_rect(slide, tbl_left, row_top, tbl_w, row_h, RH_RED_LIGHT, RH_RED)
tb = add_textbox(slide, tbl_left + Pt(10), row_top + Pt(6), Inches(1.2), Pt(20))
set_text(tb.text_frame, "Class I", FONT_TEXT, Pt(12), RH_RED, bold=True)
tb = add_textbox(slide, tbl_left + Inches(1.4), row_top + Pt(6), Inches(2), Pt(20))
set_text(tb.text_frame, "Open Science", FONT_TEXT, Pt(12), RH_RED, bold=True)
tb = add_textbox(slide, tbl_left + Inches(3.5), row_top + Pt(6), Inches(5), Pt(20))
set_text(tb.text_frame, "Fully open: data, code, parameters, research, documentation", FONT_TEXT, Pt(12), RH_RED)

# Callout
add_rect(slide, Inches(0.6), Inches(4.0), Pt(4), Inches(0.6), RH_RED)
shape = add_rect(slide, Inches(0.65), Inches(4.0), Inches(8.8), Inches(0.6), RGBColor(0xF5, 0xF5, 0xF5))
tb = add_textbox(slide, Inches(0.9), Inches(4.05), Inches(8.4), Inches(0.5))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "This enables financial institutions to select models based on transparency, auditability, and compliance readiness \u2014 not just benchmark performance.", FONT_TEXT, Pt(12), RH_BODY, bold=True)

slide.notes_slide.notes_text_frame.text = (
    "KEY MESSAGE: The Model Openness Framework addresses 'open-washing' — models that claim openness but hide critical information.\n\n"
    "- This is a real compliance risk: banks adopting third-party models may not know what data they were trained on or what licensing restrictions apply\n"
    "- Class III (Open Model): Only the model weights and basic documentation are shared — training data and code remain hidden\n"
    "- Class II (Open Tooling): Training code, evaluation code, and inference tools are shared — but training data may still be withheld\n"
    "- Class I (Open Science): Everything is open — data, code, parameters, research papers, and documentation — fully reproducible and auditable\n"
    "- For regulated financial institutions, Class I models provide the highest level of auditability for regulators\n"
    "- The MOF is integrated directly into the Open FinLLM Leaderboard — models are filtered by openness class alongside performance\n"
    "- This gives procurement and risk teams a framework for vendor selection based on transparency, not just capability"
)


# ============================================================
# SLIDE 13: Governance Stack
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_red_accent_bar(slide)
add_section_label(slide, "Governance")
add_slide_title(slide, "The Financial AI Governance Stack")
add_slide_subtitle(slide, "A layered architecture for safe deployment of AI-powered financial systems.")

# Stack layers (top to bottom)
stack_left = Inches(1.2)
stack_w = Inches(7.5)
layer_h = Inches(0.8)
gap = Pt(3)
stack_top = Inches(1.7)

layers = [
    ("BUSINESS APPLICATIONS", "Financial AI Use Cases",
     "Trading assistants  |  Risk analysis  |  Research agents  |  Advisory systems",
     RH_BLUE_LIGHT, RH_BLUE),
    ("AGENT GOVERNANCE \u2014 AgentOps", "Operational Controls for AI Agents",
     "Decision tracing  |  Tool monitoring  |  LLM-as-a-Judge  |  Human escalation",
     RH_RED_LIGHT, RH_RED),
    ("MODEL GOVERNANCE \u2014 FinLLM Benchmarking", "Ensuring Models Are Safe to Deploy",
     "Financial benchmarks  |  Numerical reasoning  |  Hallucination testing",
     RH_ORANGE_LIGHT, RH_ORANGE),
    ("TRANSPARENCY & COMPLIANCE \u2014 Model Openness Framework", "Model Provenance and Auditability",
     "Training data transparency  |  Licensing clarity  |  Reproducibility",
     RH_GREEN_LIGHT, RH_GREEN),
]

for i, (label, title, desc, bg, border) in enumerate(layers):
    y = stack_top + i * (layer_h + gap)
    shape = add_rounded_rect(slide, stack_left, y, stack_w, layer_h, bg, border)

    # Label
    tb = add_textbox(slide, stack_left + Pt(14), y + Pt(6), stack_w - Pt(28), Pt(14))
    set_text(tb.text_frame, label, FONT_MONO, Pt(8), border, bold=True)

    # Title
    tb = add_textbox(slide, stack_left + Pt(14), y + Pt(20), stack_w - Pt(28), Pt(20))
    set_text(tb.text_frame, title, FONT_DISPLAY, Pt(13), RH_BLACK, bold=True)

    # Description
    tb = add_textbox(slide, stack_left + Pt(14), y + Pt(40), stack_w - Pt(28), Pt(18))
    set_text(tb.text_frame, desc, FONT_TEXT, Pt(11), RH_BODY)

slide.notes_slide.notes_text_frame.text = (
    "KEY MESSAGE: The governance stack shows how all the frameworks fit together as a layered architecture.\n\n"
    "- Read this diagram from bottom to top — each layer depends on the one below it\n"
    "- Bottom layer (Transparency): Before deploying any model, verify its provenance — training data, licensing, reproducibility\n"
    "- Model Governance layer: Once a model passes transparency checks, benchmark it against financial tasks — test for hallucinations and numerical accuracy\n"
    "- Agent Governance layer: When models are deployed as agents, add operational controls — decision tracing, safety scoring, human escalation\n"
    "- Top layer (Applications): Only then can you safely deploy trading assistants, risk tools, research agents, and advisory systems\n"
    "- This is analogous to how banks layer infrastructure: network, platform, application, with controls at each level\n"
    "- The key takeaway: you need all four layers — skipping any one creates an unacceptable risk gap"
)


# ============================================================
# SLIDE 14: Section — What This Means
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(slide, 0, 0, Inches(4), SH, RH_RED)
tb = add_textbox(slide, Inches(1.2), Inches(1.5), Inches(2), Inches(1.5))
set_text(tb.text_frame, "06", FONT_DISPLAY, Pt(72), RGBColor(0xFF, 0x66, 0x66), bold=True, alignment=PP_ALIGN.RIGHT)
tb = add_textbox(slide, Inches(4.5), Inches(1.5), Inches(5), Inches(1.8))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "What This Means for Banks", FONT_DISPLAY, Pt(28), RH_BLACK, bold=True)
add_paragraph(tf, "Practical governance layers for financial AI deployment", FONT_TEXT, Pt(14), RH_DARK, space_before=Pt(12))

slide.notes_slide.notes_text_frame.text = (
    "KEY MESSAGE: Transition to practical recommendations for financial institutions.\n\n"
    "- We have covered the problem, the lifecycle, benchmarking, agents, and governance frameworks\n"
    "- Now we translate that into actionable guidance for a bank deploying GenAI\n"
    "- The next slide distils everything into three practical governance layers\n"
    "- These align with existing bank functions: model validation, operational risk, and vendor risk management\n"
    "- This is the 'so what?' for a CIO, risk committee, or architecture review board"
)


# ============================================================
# SLIDE 15: Three Governance Layers
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_red_accent_bar(slide)
add_section_label(slide, "For Financial Institutions")
add_slide_title(slide, "Three governance layers for financial AI")
add_slide_subtitle(slide, "Banks deploying GenAI should implement three governance layers.")

# Summary cards
sc_left = Inches(0.6)
sc_w = Inches(8.8)
sc_h = Inches(0.9)
sc_top = Inches(1.8)

items = [
    ("1", "Model Evaluation",
     "Benchmark models against financial datasets. Can the model reason about financial statements? Does it hallucinate in regulatory documents?",
     RH_RED),
    ("2", "Agent Monitoring",
     "Track AI agents operating in production. Decision tracing, audit logs, escalation procedures, and runtime evaluation of agent behaviour.",
     RH_ORANGE),
    ("3", "Model Transparency",
     "Deploy only models meeting transparency requirements. Training data provenance, licensing clarity, reproducibility, and regulatory auditability.",
     RH_GREEN),
]

for i, (num, title, desc, color) in enumerate(items):
    y = sc_top + i * (sc_h + Pt(8))
    shape = add_rounded_rect(slide, sc_left, y, sc_w, sc_h, RH_WHITE, color)

    # Number
    tb = add_textbox(slide, sc_left + Pt(14), y + Pt(14), Pt(30), Pt(30))
    set_text(tb.text_frame, num, FONT_DISPLAY, Pt(22), color, bold=True)

    # Title + desc
    tb = add_textbox(slide, sc_left + Pt(50), y + Pt(8), sc_w - Pt(70), sc_h - Pt(16))
    tf = tb.text_frame
    tf.word_wrap = True
    set_text(tf, title, FONT_DISPLAY, Pt(14), RH_BLACK, bold=True)
    add_paragraph(tf, desc, FONT_TEXT, Pt(11), RH_BODY, space_before=Pt(4))

slide.notes_slide.notes_text_frame.text = (
    "KEY MESSAGE: Banks deploying GenAI should implement three governance layers — model evaluation, agent monitoring, and model transparency.\n\n"
    "- Layer 1 — Model Evaluation: Before deploying any model, benchmark it on financial tasks. Key questions: Can it reason about 10-K filings? Does it hallucinate when analysing earnings reports? How does its numerical accuracy compare to alternatives?\n"
    "- Layer 2 — Agent Monitoring: Once agents are in production, you need continuous monitoring. This includes decision tracing, audit logs, escalation procedures, and runtime behaviour evaluation\n"
    "- Layer 3 — Model Transparency: Only deploy models where you understand the training data, licensing terms, and can demonstrate reproducibility to regulators\n"
    "- These three layers map to existing bank governance: model validation teams, operational risk teams, and vendor risk management\n"
    "- The investment required is in tooling and process — the frameworks described in this presentation are open source and available today\n"
    "- Start with Layer 3 (transparency) for procurement decisions, then build Layers 1 and 2 as you scale deployment"
)


# ============================================================
# SLIDE 16: Core Message (dark background)
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(slide, 0, 0, SW, SH, RH_BLACK)

# Title
tb = add_textbox(slide, Inches(0.8), Inches(0.6), Inches(8.4), Inches(1.0))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "Financial AI should be treated like\nregulated financial infrastructure", FONT_DISPLAY, Pt(26), RH_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Four pillars
pillar_w = Inches(2.1)
pillar_h = Inches(1.4)
pillar_top = Inches(2.0)
pillar_gap = Inches(0.15)
pillar_left = Inches(0.55)

pillars = [
    ("Benchmarked", "Against standardised financial tasks"),
    ("Monitored", "Through operational evaluation in production"),
    ("Governed", "Through compliance and transparency frameworks"),
    ("Auditable", "With traceable model behaviour"),
]

for i, (title, desc) in enumerate(pillars):
    x = pillar_left + i * (pillar_w + pillar_gap)
    shape = add_rounded_rect(slide, x, pillar_top, pillar_w, pillar_h,
                             RGBColor(0x25, 0x25, 0x25), RGBColor(0x44, 0x44, 0x44))

    tb = add_textbox(slide, x + Pt(12), pillar_top + Pt(16), pillar_w - Pt(24), pillar_h - Pt(32))
    tf = tb.text_frame
    tf.word_wrap = True
    set_text(tf, title, FONT_DISPLAY, Pt(16), RH_RED, bold=True, alignment=PP_ALIGN.CENTER)
    add_paragraph(tf, desc, FONT_TEXT, Pt(11), RGBColor(0xBB, 0xBB, 0xBB), alignment=PP_ALIGN.CENTER, space_before=Pt(8))

# Tagline
tb = add_textbox(slide, Inches(1), Inches(3.8), Inches(8), Inches(0.6))
set_text(tb.text_frame, "This is not consumer AI.", FONT_DISPLAY, Pt(20), RH_RED, bold=True, alignment=PP_ALIGN.CENTER)

tb = add_textbox(slide, Inches(1), Inches(4.3), Inches(8), Inches(0.4))
set_text(tb.text_frame, "This is financial infrastructure.", FONT_TEXT, Pt(16), RGBColor(0x99, 0x99, 0x99), alignment=PP_ALIGN.CENTER)

slide.notes_slide.notes_text_frame.text = (
    "KEY MESSAGE: Financial AI must be treated as regulated financial infrastructure — benchmarked, monitored, governed, and auditable.\n\n"
    "- This is the single most important slide — pause here and let the message land\n"
    "- Benchmarked: Models must be tested against standardised financial tasks before deployment — just as we stress-test risk models\n"
    "- Monitored: AI agents in production require continuous operational evaluation — just as we monitor trading systems\n"
    "- Governed: Compliance and transparency frameworks must be applied — just as we govern financial data and processes\n"
    "- Auditable: Every model decision must be traceable — regulators will require this\n"
    "- The closing line is deliberate: 'This is not consumer AI. This is financial infrastructure.' — frame AI governance as a business requirement, not a technology experiment"
)


# ============================================================
# SLIDE 17: Thank You (red splash)
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(slide, 0, 0, SW, SH, RH_RED)

tb = add_textbox(slide, Inches(1), Inches(1.5), Inches(8), Inches(1.2))
set_text(tb.text_frame, "Thank You", FONT_DISPLAY, Pt(44), RH_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_rect(slide, Inches(4.2), Inches(2.7), Inches(1.6), Pt(3), RGBColor(0xFF, 0x99, 0x99))

tb = add_textbox(slide, Inches(1.5), Inches(3.0), Inches(7), Inches(0.6))
set_text(tb.text_frame, "Questions & Discussion", FONT_TEXT, Pt(20), RH_WHITE, alignment=PP_ALIGN.CENTER)

slide.notes_slide.notes_text_frame.text = (
    "KEY MESSAGE: Close the presentation and open discussion.\n\n"
    "- Recap the three governance layers: model evaluation, agent monitoring, model transparency\n"
    "- Suggested discussion questions: Where is your organisation on the Exploration/Readiness/Governance maturity model?\n"
    "- What AI models are currently deployed or under evaluation? Have they been benchmarked on financial tasks?\n"
    "- Do you have operational monitoring for AI agents? Can you trace agent decisions for audit purposes?\n"
    "- Are your AI model procurement decisions considering training data provenance and licensing?\n"
    "- The frameworks discussed are open source — the Open FinLLM Leaderboard, AgentOps, and Model Openness Framework are all available today\n"
    "- Offer to share the original paper and links to the leaderboard and governance framework resources"
)


# ============================================================
# SAVE
# ============================================================
prs.save(OUTPUT)
print(f"Presentation saved to: {OUTPUT}")
print(f"Total slides: {len(prs.slides)}")
